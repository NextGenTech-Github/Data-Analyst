                                                       from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

# Generate time series data
x = np.linspace(0, 10, 100)
y = np.sin(x)

# Create a plot of the time series data
plt.figure()
plt.plot(x, y)
plt.title('Time Series Data')
plt.xlabel('Time')
plt.ylabel('Value')

# Save the plot as an image
plot_path = 'time_series_plot.png'
plt.savefig(plot_path)

# Create a PowerPoint presentation
prs = Presentation()

# Add a slide
slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
slide = prs.slides.add_slide(slide_layout)

# Add the plot image to the slide
left = Inches(1)
top = Inches(1)
height = Inches(5)
pic = slide.shapes.add_picture(plot_path, left, top, height=height)

# Save the presentation
pptx_path = 'Presentation.pptx'
prs.save(pptx_path)

pptx_path
