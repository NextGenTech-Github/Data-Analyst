from pptx import Presentation
from pptx.util import Inches
import matplotlib.pyplot as plt
import numpy as np

# Create a PowerPoint presentation
prs = Presentation()

# Function to create and save a plot
def create_plot(x, y, title, xlabel, ylabel, file_path):
    plt.figure()
    plt.plot(x, y)
    plt.title(title)
    plt.xlabel(xlabel)
    plt.ylabel(ylabel)
    plt.savefig(file_path)
    plt.close()

# Time series data
x = np.linspace(0, 10, 100)

# Data for multiple graphs
data = [
    (x, np.sin(x), 'Sine Wave', 'Time', 'Value', 'sine_wave.png'),
    (x, np.cos(x), 'Cosine Wave', 'Time', 'Value', '/cosine_wave.png'),
    (x, np.tan(x), 'Tangent Wave', 'Time', 'Value', 'tangent_wave.png')
]

# Generate and add plots to slides
for x, y, title, xlabel, ylabel, plot_path in data:
    create_plot(x, y, title, xlabel, ylabel, plot_path)
    
    slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(1)
    height = Inches(5)
    pic = slide.shapes.add_picture(plot_path, left, top, height=height)

# Save the presentation
pptx_path = 'multiple_graphs_presentation.pptx'
prs.save(pptx_path)

pptx_path
